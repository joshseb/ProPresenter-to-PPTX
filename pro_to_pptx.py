"""
ProPresenter 7 to PowerPoint Converter
Converts .pro files (ProPresenter 7) to .pptx (Microsoft PowerPoint).
"""

import os
import sys
import time
import json
import threading
import tkinter as tk
from tkinter import filedialog, messagebox, ttk

# Add proto_pb2 folder to path so compiled protobuf bindings are importable
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "proto_pb2"))

try:
    import presentation_pb2
    import action_pb2
    import presentationSlide_pb2
    PROTO_AVAILABLE = True
except ImportError:
    PROTO_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    STRIPRTF_AVAILABLE = True
except ImportError:
    STRIPRTF_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from watchdog.observers import Observer
    from watchdog.events import FileSystemEventHandler
    WATCHDOG_AVAILABLE = True
except ImportError:
    WATCHDOG_AVAILABLE = False

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SLIDE_WIDTH_IN  = 13.33   # 16:9 widescreen
SLIDE_HEIGHT_IN = 7.5
PRO_CANVAS_W    = 1920.0
PRO_CANVAS_H    = 1080.0
DEFAULT_FONT_PT = 60
FONT_COLOR      = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_AVAILABLE else None
BG_COLOR        = RGBColor(0x00, 0x00, 0x00) if PPTX_AVAILABLE else None

DEFAULT_WATCH_FOLDER = os.path.expanduser("~/Documents/ProPresenter")

# History file lives alongside pro_to_pptx.py
HISTORY_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "converted_history.json")

# ---------------------------------------------------------------------------
# ConversionHistory – tracks which files have already been converted
# ---------------------------------------------------------------------------

class ConversionHistory:
    """
    Persists a JSON file mapping absolute .pro path → last-converted mtime.
    A file is considered "already converted" if:
      1. Its path is in the history, AND
      2. Its current mtime matches the recorded mtime (i.e. it hasn't changed).
    If the .pro file was modified since it was last converted, it is treated
    as a new file and will be converted again.
    """

    def __init__(self, history_path: str = HISTORY_FILE):
        self._path = history_path
        self._lock = threading.Lock()
        self._data: dict = {}   # { absolute_path: mtime_float }
        self._load()

    def _load(self):
        if os.path.isfile(self._path):
            try:
                with open(self._path, "r", encoding="utf-8") as f:
                    self._data = json.load(f)
            except Exception:
                self._data = {}

    def _save(self):
        try:
            with open(self._path, "w", encoding="utf-8") as f:
                json.dump(self._data, f, indent=2)
        except Exception:
            pass  # non-fatal; history is a convenience, not critical

    def is_converted(self, pro_path: str) -> bool:
        """Return True if this file has already been successfully converted and hasn't changed."""
        abs_path = os.path.abspath(pro_path)
        with self._lock:
            recorded_mtime = self._data.get(abs_path)
        if recorded_mtime is None:
            return False
        try:
            current_mtime = os.path.getmtime(pro_path)
        except OSError:
            return False
        return abs(current_mtime - recorded_mtime) < 1.0  # within 1 second = same version

    def mark_converted(self, pro_path: str):
        """Record that this file was successfully converted right now."""
        abs_path = os.path.abspath(pro_path)
        try:
            mtime = os.path.getmtime(pro_path)
        except OSError:
            mtime = time.time()
        with self._lock:
            self._data[abs_path] = mtime
            self._save()

    def clear(self):
        """Wipe the entire history."""
        with self._lock:
            self._data = {}
            self._save()

    def count(self) -> int:
        with self._lock:
            return len(self._data)


# Shared singleton used by both single-file and watcher paths
_history = ConversionHistory()


# ---------------------------------------------------------------------------
# ProParser – reads a .pro file and extracts per-slide text blocks
# ---------------------------------------------------------------------------

class ProParser:

    def parse(self, filepath: str) -> list:
        """
        Returns a list of slides.
        Each slide is a list of dicts:
          { "text": str, "left_frac": float, "top_frac": float,
            "width_frac": float, "height_frac": float }
        """
        if not PROTO_AVAILABLE:
            raise RuntimeError(
                "Proto bindings not found in proto_pb2/.\n"
                "Run the setup step to compile the .proto files first.\n"
                "See README or the instructions in the app."
            )

        with open(filepath, "rb") as f:
            raw = f.read()

        pres = presentation_pb2.Presentation()
        try:
            pres.ParseFromString(raw)
        except Exception as exc:
            raise RuntimeError(
                f"Could not parse '{os.path.basename(filepath)}' as a "
                f"ProPresenter 7 file.\nDetails: {exc}"
            ) from exc

        slides = []
        for cue in pres.cues:
            for action in cue.actions:
                if action.type != action_pb2.Action.ActionType.Value(
                        "ACTION_TYPE_PRESENTATION_SLIDE"):
                    continue
                slide_data = self._extract_slide(
                    action.slide.presentation.base_slide)
                if slide_data:
                    slides.append(slide_data)

        return slides

    def _extract_slide(self, base_slide) -> list:
        text_blocks = []
        for elem_wrapper in base_slide.elements:
            graphics_elem = elem_wrapper.element

            # Skip elements with no text
            rtf_bytes = graphics_elem.text.rtf_data
            if not rtf_bytes:
                continue

            plain = self._rtf_to_plain(rtf_bytes)
            if not plain.strip():
                continue

            bounds = graphics_elem.bounds
            left_frac   = bounds.origin.x      / PRO_CANVAS_W
            top_frac    = bounds.origin.y      / PRO_CANVAS_H
            width_frac  = bounds.size.width    / PRO_CANVAS_W
            height_frac = bounds.size.height   / PRO_CANVAS_H

            # Clamp fractions to [0, 1] in case of out-of-bounds elements
            left_frac   = max(0.0, min(left_frac,  1.0))
            top_frac    = max(0.0, min(top_frac,   1.0))
            width_frac  = max(0.05, min(width_frac, 1.0 - left_frac))
            height_frac = max(0.05, min(height_frac, 1.0 - top_frac))

            text_blocks.append({
                "text":        plain.strip(),
                "left_frac":   left_frac,
                "top_frac":    top_frac,
                "width_frac":  width_frac,
                "height_frac": height_frac,
            })
        return text_blocks

    @staticmethod
    def _rtf_to_plain(rtf_bytes: bytes) -> str:
        if not STRIPRTF_AVAILABLE:
            # Fallback: strip RTF tags with a simple regex
            import re
            try:
                text = rtf_bytes.decode("utf-8", errors="ignore")
            except Exception:
                text = rtf_bytes.decode("latin-1", errors="ignore")
            text = re.sub(r"\\[a-z]+\d*\s?", " ", text)
            text = re.sub(r"[{}\\]", "", text)
            return text.strip()

        try:
            rtf_str = rtf_bytes.decode("utf-8", errors="ignore")
        except Exception:
            rtf_str = rtf_bytes.decode("latin-1", errors="ignore")
        return rtf_to_text(rtf_str, errors="ignore")


# ---------------------------------------------------------------------------
# PptxBuilder – writes slides to a .pptx file
# ---------------------------------------------------------------------------

class PptxBuilder:

    def build(self, slides: list, output_path: str) -> int:
        if not PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx is not installed.\nRun: pip install python-pptx"
            )

        prs = PptxPresentation()
        prs.slide_width  = Inches(SLIDE_WIDTH_IN)
        prs.slide_height = Inches(SLIDE_HEIGHT_IN)

        # Find the blank layout (index 6 by convention; fall back to search)
        blank_layout = next(
            (l for l in prs.slide_layouts if l.name == "Blank"),
            prs.slide_layouts[6]
        )

        for slide_data in slides:
            slide = prs.slides.add_slide(blank_layout)
            self._set_background(slide)
            for block in slide_data:
                self._add_textbox(slide, block)

        prs.save(output_path)
        return len(slides)

    @staticmethod
    def _set_background(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    @staticmethod
    def _add_textbox(slide, block: dict):
        left   = Inches(SLIDE_WIDTH_IN  * block["left_frac"])
        top    = Inches(SLIDE_HEIGHT_IN * block["top_frac"])
        width  = Inches(SLIDE_WIDTH_IN  * block["width_frac"])
        height = Inches(SLIDE_HEIGHT_IN * block["height_frac"])

        txBox = slide.shapes.add_textbox(left, top, width, height)

        # Make the text box fully transparent — no fill, no border line —
        # so only the white text is visible over the black slide background.
        from pptx.oxml.ns import qn
        from lxml import etree

        sp = txBox._element
        # Remove any existing spPr fill/line children and set no-fill + no-line
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            spPr = etree.SubElement(sp, qn("p:spPr"))

        # No fill: <a:noFill/>
        for tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"),
                    qn("a:pattFill"), qn("a:blipFill"), qn("a:ln")):
            for old in spPr.findall(tag):
                spPr.remove(old)

        no_fill = etree.SubElement(spPr, qn("a:noFill"))

        # No border line: <a:ln><a:noFill/></a:ln>
        ln = etree.SubElement(spPr, qn("a:ln"))
        etree.SubElement(ln, qn("a:noFill"))

        tf = txBox.text_frame
        tf.word_wrap = True

        lines = block["text"].split("\n")
        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = line
            run.font.size = Pt(DEFAULT_FONT_PT)
            run.font.color.rgb = FONT_COLOR
            run.font.bold = False


# ---------------------------------------------------------------------------
# Conversion helper – used by both single-file and folder-watcher paths
# ---------------------------------------------------------------------------

def convert_file(pro_path: str, output_dir: str, watch_root: str = None) -> tuple:
    """
    Convert one .pro file to .pptx in output_dir (always flat — no subfolders).
    watch_root is accepted but ignored (kept for call-site compatibility).
    Returns (output_path, slide_count) on success.
    Raises RuntimeError on failure.
    On success, records the file in the conversion history.
    """
    base = os.path.splitext(os.path.basename(pro_path))[0]
    os.makedirs(output_dir, exist_ok=True)
    out_path = os.path.join(output_dir, base + ".pptx")

    parser = ProParser()
    slides = parser.parse(pro_path)

    builder = PptxBuilder()
    count = builder.build(slides, out_path)

    # Record successful conversion so it won't be re-converted automatically
    _history.mark_converted(pro_path)

    return out_path, count


# ---------------------------------------------------------------------------
# FolderWatcher – monitors a folder for new .pro files
# ---------------------------------------------------------------------------

class _ProFileHandler(FileSystemEventHandler):

    def __init__(self, watch_root: str, output_dir: str, callback):
        super().__init__()
        self._watch_root = watch_root
        self._output_dir = output_dir
        self._callback = callback  # callback(pro_path, result_or_error)
        self._pending: dict = {}   # path -> scheduled time
        self._lock = threading.Lock()

    def on_created(self, event):
        if not event.is_directory and event.src_path.endswith(".pro"):
            self._schedule(event.src_path)

    def on_moved(self, event):
        if not event.is_directory and event.dest_path.endswith(".pro"):
            self._schedule(event.dest_path)

    def _schedule(self, path: str):
        """Wait 1 second after the last event before processing (file may still be writing)."""
        with self._lock:
            self._pending[path] = time.time() + 1.0
        threading.Thread(target=self._delayed_convert, args=(path,), daemon=True).start()

    def _delayed_convert(self, path: str):
        time.sleep(1.1)
        with self._lock:
            due = self._pending.get(path)
        if due is None or time.time() < due:
            return  # superseded or not yet due
        with self._lock:
            self._pending.pop(path, None)

        # Skip files that were already converted and haven't changed
        if _history.is_converted(path):
            self._callback(path, "SKIPPED (already converted — file unchanged)")
            return

        try:
            out_path, count = convert_file(path, self._output_dir, self._watch_root)
            self._callback(path, f"OK: {count} slides → {os.path.basename(out_path)}")
        except Exception as exc:
            self._callback(path, f"ERROR: {exc}")


class FolderWatcher:

    def __init__(self, watch_dir: str, output_dir: str, callback):
        self._watch_dir  = watch_dir
        self._output_dir = output_dir
        self._callback   = callback
        self._observer: Observer = None

    def start(self):
        if self._observer and self._observer.is_alive():
            return
        handler = _ProFileHandler(self._watch_dir, self._output_dir, self._callback)
        self._observer = Observer()
        self._observer.schedule(handler, self._watch_dir, recursive=True)
        self._observer.start()

    def stop(self):
        if self._observer:
            self._observer.stop()
            self._observer.join(timeout=3)
            self._observer = None

    @property
    def running(self) -> bool:
        return self._observer is not None and self._observer.is_alive()


# ---------------------------------------------------------------------------
# merge_pptx_files – combines multiple .pptx files into one
# ---------------------------------------------------------------------------

def merge_pptx_files(input_paths: list, output_path: str) -> int:
    """
    Combine all slides from input_paths (in order) into a single .pptx.
    Returns the total number of slides written.
    Uses python-pptx's XML copy approach to preserve slide content faithfully.
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed.\nRun: pip install python-pptx")

    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches
    from lxml import etree
    import copy

    # Open the first file as the base presentation
    merged = PptxPresentation(input_paths[0])
    total_slides = len(merged.slides)

    for src_path in input_paths[1:]:
        src_prs = PptxPresentation(src_path)
        for src_slide in src_prs.slides:
            # Add a blank slide to the merged presentation
            blank_layout = next(
                (l for l in merged.slide_layouts if l.name == "Blank"),
                merged.slide_layouts[6]
            )
            new_slide = merged.slides.add_slide(blank_layout)

            # Copy slide XML content from source to new slide
            # Remove all existing shapes from the blank slide
            sp_tree = new_slide.shapes._spTree
            for child in list(sp_tree):
                sp_tree.remove(child)

            # Copy all elements from source slide's shape tree
            src_sp_tree = src_slide.shapes._spTree
            for elem in src_sp_tree:
                sp_tree.append(copy.deepcopy(elem))

            # Copy slide background if set
            src_bg = src_slide.background
            new_bg = new_slide.background
            src_fill = src_bg.fill
            if src_fill.type is not None:
                new_bg_elem = new_bg._element
                src_bg_elem = src_bg._element
                # Replace background element content
                for child in list(new_bg_elem):
                    new_bg_elem.remove(child)
                for child in src_bg_elem:
                    new_bg_elem.append(copy.deepcopy(child))

            total_slides += 1

    merged.save(output_path)
    return total_slides


# ---------------------------------------------------------------------------
# ConverterApp – tkinter GUI
# ---------------------------------------------------------------------------

class ConverterApp:

    def __init__(self, root: tk.Tk):
        self.root = root
        self.root.title("ProPresenter 7  →  PowerPoint Converter")
        self.root.geometry("640x430")
        self.root.resizable(False, False)

        self._watcher: FolderWatcher = None
        self._build_ui()
        self._check_deps()

    # ------------------------------------------------------------------
    # UI construction
    # ------------------------------------------------------------------

    def _build_ui(self):
        notebook = ttk.Notebook(self.root)
        notebook.pack(fill="both", expand=True, padx=10, pady=(10, 0))

        # Tab 1 – Single file
        tab1 = tk.Frame(notebook)
        notebook.add(tab1, text="  Single File  ")
        self._build_single_tab(tab1)

        # Tab 2 – Watch folder
        tab2 = tk.Frame(notebook)
        notebook.add(tab2, text="  Watch Folder  ")
        self._build_watch_tab(tab2)

        # Tab 3 – Merge PPTX files
        tab3 = tk.Frame(notebook)
        notebook.add(tab3, text="  Merge PPTX  ")
        self._build_merge_tab(tab3)

        # Shared bottom bar
        bottom = tk.Frame(self.root)
        bottom.pack(fill="x", padx=10, pady=6)

        self.progress = ttk.Progressbar(bottom, mode="indeterminate", length=620)
        self.progress.pack(fill="x")

        self.status_var = tk.StringVar(value="Ready.")
        tk.Label(bottom, textvariable=self.status_var, anchor="w",
                 wraplength=620, justify="left", fg="#444").pack(fill="x", pady=(4, 0))

    def _build_single_tab(self, parent):
        pad = {"padx": 10, "pady": 8}

        tk.Label(parent, text="ProPresenter File (.pro):", anchor="e").grid(
            row=0, column=0, sticky="e", **pad)
        self.pro_path = tk.StringVar()
        ttk.Entry(parent, textvariable=self.pro_path, width=46).grid(
            row=0, column=1, sticky="ew", padx=(0, 4), pady=8)
        ttk.Button(parent, text="Browse…", command=self._pick_pro_file).grid(
            row=0, column=2, **pad)

        tk.Label(parent, text="Output Folder:", anchor="e").grid(
            row=1, column=0, sticky="e", **pad)
        self.single_out_dir = tk.StringVar()
        ttk.Entry(parent, textvariable=self.single_out_dir, width=46).grid(
            row=1, column=1, sticky="ew", padx=(0, 4), pady=8)
        ttk.Button(parent, text="Browse…", command=self._pick_single_out_dir).grid(
            row=1, column=2, **pad)

        self.convert_btn = ttk.Button(
            parent, text="Convert",
            command=self._start_single_conversion)
        self.convert_btn.grid(row=2, column=0, columnspan=3, pady=16)

        parent.columnconfigure(1, weight=1)

    def _build_watch_tab(self, parent):
        pad = {"padx": 10, "pady": 8}

        tk.Label(parent, text="Watch Folder:", anchor="e").grid(
            row=0, column=0, sticky="e", **pad)
        self.watch_dir = tk.StringVar(value=DEFAULT_WATCH_FOLDER)
        ttk.Entry(parent, textvariable=self.watch_dir, width=46).grid(
            row=0, column=1, sticky="ew", padx=(0, 4), pady=8)
        ttk.Button(parent, text="Browse…", command=self._pick_watch_dir).grid(
            row=0, column=2, **pad)

        tk.Label(parent, text="Output Folder:", anchor="e").grid(
            row=1, column=0, sticky="e", **pad)
        self.watch_out_dir = tk.StringVar()
        ttk.Entry(parent, textvariable=self.watch_out_dir, width=46).grid(
            row=1, column=1, sticky="ew", padx=(0, 4), pady=8)
        ttk.Button(parent, text="Browse…", command=self._pick_watch_out_dir).grid(
            row=1, column=2, **pad)

        btn_frame = tk.Frame(parent)
        btn_frame.grid(row=2, column=0, columnspan=3, pady=6)
        self.watch_btn = ttk.Button(
            btn_frame, text="Start Watching",
            command=self._toggle_watcher)
        self.watch_btn.pack(side="left", padx=(0, 6))

        self.convert_existing_btn = ttk.Button(
            btn_frame, text="Convert Existing",
            command=self._convert_existing)
        self.convert_existing_btn.pack(side="left", padx=(0, 6))

        self.clear_history_btn = ttk.Button(
            btn_frame, text="Clear History",
            command=self._clear_history)
        self.clear_history_btn.pack(side="left")

        tk.Label(parent, text="Activity Log:", anchor="w").grid(
            row=3, column=0, columnspan=3, sticky="w", padx=10)

        log_frame = tk.Frame(parent)
        log_frame.grid(row=4, column=0, columnspan=3, sticky="nsew", padx=10, pady=(0, 6))
        scrollbar = ttk.Scrollbar(log_frame)
        scrollbar.pack(side="right", fill="y")
        self.log_text = tk.Text(log_frame, height=5, width=72,
                                yscrollcommand=scrollbar.set,
                                state="disabled", bg="#f4f4f4", font=("Courier", 10))
        self.log_text.pack(side="left", fill="both")
        scrollbar.config(command=self.log_text.yview)

        parent.columnconfigure(1, weight=1)
        parent.rowconfigure(4, weight=1)

    def _build_merge_tab(self, parent):
        pad = {"padx": 10, "pady": 4}

        # Instruction label
        tk.Label(parent, text="Add .pptx files below, reorder them, then click Merge.",
                 anchor="w", fg="#555").grid(row=0, column=0, columnspan=3,
                                             sticky="w", padx=10, pady=(8, 2))

        # List box + scrollbar
        list_frame = tk.Frame(parent)
        list_frame.grid(row=1, column=0, columnspan=2, sticky="nsew", padx=10, pady=4)
        scrollbar = ttk.Scrollbar(list_frame)
        scrollbar.pack(side="right", fill="y")
        self.merge_listbox = tk.Listbox(
            list_frame, yscrollcommand=scrollbar.set,
            selectmode="extended", height=8, font=("Courier", 10))
        self.merge_listbox.pack(side="left", fill="both", expand=True)
        scrollbar.config(command=self.merge_listbox.yview)

        # Right-side buttons: Add / Remove / Up / Down
        side_btns = tk.Frame(parent)
        side_btns.grid(row=1, column=2, sticky="n", padx=(0, 10), pady=4)
        ttk.Button(side_btns, text="Add Files…",
                   command=self._merge_add_files).pack(pady=2, fill="x")
        ttk.Button(side_btns, text="Remove",
                   command=self._merge_remove_selected).pack(pady=2, fill="x")
        ttk.Button(side_btns, text="▲  Move Up",
                   command=self._merge_move_up).pack(pady=(10, 2), fill="x")
        ttk.Button(side_btns, text="▼  Move Down",
                   command=self._merge_move_down).pack(pady=2, fill="x")
        ttk.Button(side_btns, text="Clear All",
                   command=self._merge_clear).pack(pady=(10, 2), fill="x")

        # Output file row
        tk.Label(parent, text="Save merged file as:", anchor="e").grid(
            row=2, column=0, sticky="e", **pad)
        self.merge_out_path = tk.StringVar()
        ttk.Entry(parent, textvariable=self.merge_out_path, width=42).grid(
            row=2, column=1, sticky="ew", padx=(0, 4), pady=4)
        ttk.Button(parent, text="Browse…",
                   command=self._merge_pick_output).grid(row=2, column=2, padx=(0, 10), pady=4)

        # Merge button
        self.merge_btn = ttk.Button(
            parent, text="Merge into One PPTX",
            command=self._start_merge)
        self.merge_btn.grid(row=3, column=0, columnspan=3, pady=10)

        parent.columnconfigure(1, weight=1)
        parent.rowconfigure(1, weight=1)

    # ------------------------------------------------------------------
    # File/folder pickers
    # ------------------------------------------------------------------

    def _pick_pro_file(self):
        path = filedialog.askopenfilename(
            title="Select a ProPresenter 7 file",
            filetypes=[("ProPresenter 7 Files", "*.pro"), ("All Files", "*.*")])
        if path:
            self.pro_path.set(path)
            if not self.single_out_dir.get():
                self.single_out_dir.set(os.path.dirname(path))

    def _pick_single_out_dir(self):
        path = filedialog.askdirectory(title="Select Output Folder")
        if path:
            self.single_out_dir.set(path)

    def _pick_watch_dir(self):
        path = filedialog.askdirectory(title="Select Folder to Watch")
        if path:
            self.watch_dir.set(path)

    def _pick_watch_out_dir(self):
        path = filedialog.askdirectory(title="Select Output Folder")
        if path:
            self.watch_out_dir.set(path)

    # ------------------------------------------------------------------
    # Single-file conversion
    # ------------------------------------------------------------------

    def _start_single_conversion(self):
        pro_file = self.pro_path.get().strip()
        out_dir  = self.single_out_dir.get().strip()

        if not pro_file or not os.path.isfile(pro_file):
            messagebox.showerror("Error", "Please select a valid .pro file.")
            return
        if not out_dir:
            messagebox.showerror("Error", "Please select an output folder.")
            return
        os.makedirs(out_dir, exist_ok=True)

        # Warn if this file has already been converted and hasn't changed
        if _history.is_converted(pro_file):
            answer = messagebox.askyesno(
                "Already Converted",
                f"'{os.path.basename(pro_file)}' has already been converted "
                f"and the file hasn't changed since.\n\nConvert again anyway?")
            if not answer:
                self.status_var.set("Skipped (already converted).")
                return

        self.convert_btn.config(state="disabled")
        self.progress.start(10)
        self.status_var.set("Converting…")

        threading.Thread(
            target=self._run_single_conversion,
            args=(pro_file, out_dir),
            daemon=True
        ).start()

    def _run_single_conversion(self, pro_file: str, out_dir: str):
        try:
            out_path, count = convert_file(pro_file, out_dir)
            self.root.after(0, self._on_single_success, count, out_path)
        except Exception as exc:
            self.root.after(0, self._on_single_error, str(exc))

    def _on_single_success(self, count: int, out_path: str):
        self.progress.stop()
        self.convert_btn.config(state="normal")
        self.status_var.set(f"Done. {count} slides written to: {out_path}")
        messagebox.showinfo("Conversion Complete",
            f"{count} slides converted.\nSaved to:\n{out_path}")

    def _on_single_error(self, msg: str):
        self.progress.stop()
        self.convert_btn.config(state="normal")
        self.status_var.set(f"Error: {msg}")
        messagebox.showerror("Conversion Failed", msg)

    # ------------------------------------------------------------------
    # Folder watcher
    # ------------------------------------------------------------------

    def _toggle_watcher(self):
        if self._watcher and self._watcher.running:
            self._stop_watcher()
        else:
            self._start_watcher()

    def _start_watcher(self):
        if not WATCHDOG_AVAILABLE:
            messagebox.showerror(
                "Missing Dependency",
                "watchdog is not installed.\nRun: pip install watchdog")
            return

        watch_dir = self.watch_dir.get().strip()
        out_dir   = self.watch_out_dir.get().strip()

        if not watch_dir or not os.path.isdir(watch_dir):
            messagebox.showerror("Error",
                "Watch folder does not exist.\n"
                "Please create it or select a valid folder.")
            return
        if not out_dir:
            messagebox.showerror("Error", "Please select an output folder.")
            return
        os.makedirs(out_dir, exist_ok=True)

        self._watcher = FolderWatcher(
            watch_dir, out_dir,
            callback=lambda path, result: self.root.after(
                0, self._on_watch_event, path, result)
        )
        self._watcher.start()

        self.watch_btn.config(text="Stop Watching")
        self.status_var.set(f"Watching: {watch_dir}")
        self._log(f"Started watching: {watch_dir}")

    def _stop_watcher(self):
        if self._watcher:
            self._watcher.stop()
            self._watcher = None
        self.watch_btn.config(text="Start Watching")
        self.status_var.set("Watcher stopped.")
        self._log("Watcher stopped.")

    def _convert_existing(self):
        """Scan the watch folder now and convert all .pro files not yet converted."""
        watch_dir = self.watch_dir.get().strip()
        out_dir   = self.watch_out_dir.get().strip()

        if not watch_dir or not os.path.isdir(watch_dir):
            messagebox.showerror("Error",
                "Watch folder does not exist.\nPlease select a valid folder first.")
            return
        if not out_dir:
            messagebox.showerror("Error", "Please select an output folder first.")
            return
        os.makedirs(out_dir, exist_ok=True)

        # Collect all .pro files recursively
        pro_files = []
        for dirpath, _, filenames in os.walk(watch_dir):
            for fn in filenames:
                if fn.endswith(".pro"):
                    pro_files.append(os.path.join(dirpath, fn))

        if not pro_files:
            messagebox.showinfo("No Files Found",
                f"No .pro files found in:\n{watch_dir}")
            return

        # Ask how many are new (skip already-converted)
        new_files = [p for p in pro_files if not _history.is_converted(p)]
        skipped   = len(pro_files) - len(new_files)

        if not new_files:
            messagebox.showinfo("Nothing to Convert",
                f"All {len(pro_files)} .pro file(s) have already been converted.\n"
                f"Use 'Clear History' to re-convert them.")
            return

        answer = messagebox.askyesno(
            "Convert Existing Files",
            f"Found {len(new_files)} unconverted .pro file(s)"
            + (f" ({skipped} already converted, skipped)" if skipped else "")
            + f" in:\n{watch_dir}\n\nConvert them now?")
        if not answer:
            return

        self.convert_existing_btn.config(state="disabled")
        self.progress.start(10)
        self.status_var.set(f"Converting {len(new_files)} existing file(s)…")

        threading.Thread(
            target=self._run_convert_existing,
            args=(new_files, out_dir, watch_dir),
            daemon=True
        ).start()

    def _run_convert_existing(self, pro_files: list, out_dir: str, watch_root: str):
        done, errors = 0, 0
        for path in pro_files:
            try:
                rel = os.path.relpath(path, watch_root)
            except ValueError:
                rel = os.path.basename(path)
            try:
                out_path, count = convert_file(path, out_dir, watch_root)
                done += 1
                self.root.after(0, self._log,
                    f"{rel}: OK — {count} slides → {os.path.basename(out_path)}")
            except Exception as exc:
                errors += 1
                self.root.after(0, self._log, f"{rel}: ERROR — {exc}")
        self.root.after(0, self._on_convert_existing_done, done, errors)

    def _on_convert_existing_done(self, done: int, errors: int):
        self.progress.stop()
        self.convert_existing_btn.config(state="normal")
        msg = f"Batch complete: {done} converted"
        if errors:
            msg += f", {errors} failed"
        self.status_var.set(msg + ".")
        self._log(msg + ".")

    def _on_watch_event(self, pro_path: str, result: str):
        watch_root = self.watch_dir.get().strip()
        try:
            rel = os.path.relpath(pro_path, watch_root)
        except ValueError:
            rel = os.path.basename(pro_path)
        self._log(f"{rel}: {result}")
        self.status_var.set(f"Last: {rel} — {result}")

    def _clear_history(self):
        count = _history.count()
        if count == 0:
            messagebox.showinfo("Clear History", "History is already empty.")
            return
        answer = messagebox.askyesno(
            "Clear History",
            f"This will forget that {count} file(s) were already converted.\n"
            f"They will be converted again next time they are detected.\n\n"
            f"Are you sure?")
        if answer:
            _history.clear()
            self._log(f"History cleared ({count} entries removed).")
            self.status_var.set("Conversion history cleared.")

    def _log(self, message: str):
        self.log_text.config(state="normal")
        self.log_text.insert("end", message + "\n")
        self.log_text.see("end")
        self.log_text.config(state="disabled")

    # ------------------------------------------------------------------
    # Merge PPTX helpers
    # ------------------------------------------------------------------

    def _merge_add_files(self):
        paths = filedialog.askopenfilenames(
            title="Select PowerPoint files to merge",
            filetypes=[("PowerPoint Files", "*.pptx"), ("All Files", "*.*")])
        for p in paths:
            # Avoid duplicates
            existing = list(self.merge_listbox.get(0, "end"))
            if p not in existing:
                self.merge_listbox.insert("end", p)

    def _merge_remove_selected(self):
        selected = list(self.merge_listbox.curselection())
        for idx in reversed(selected):
            self.merge_listbox.delete(idx)

    def _merge_move_up(self):
        selected = list(self.merge_listbox.curselection())
        if not selected or selected[0] == 0:
            return
        for idx in selected:
            text = self.merge_listbox.get(idx)
            self.merge_listbox.delete(idx)
            self.merge_listbox.insert(idx - 1, text)
            self.merge_listbox.selection_set(idx - 1)

    def _merge_move_down(self):
        selected = list(self.merge_listbox.curselection())
        if not selected or selected[-1] == self.merge_listbox.size() - 1:
            return
        for idx in reversed(selected):
            text = self.merge_listbox.get(idx)
            self.merge_listbox.delete(idx)
            self.merge_listbox.insert(idx + 1, text)
            self.merge_listbox.selection_set(idx + 1)

    def _merge_clear(self):
        self.merge_listbox.delete(0, "end")

    def _merge_pick_output(self):
        path = filedialog.asksaveasfilename(
            title="Save merged PPTX as…",
            defaultextension=".pptx",
            filetypes=[("PowerPoint Files", "*.pptx")])
        if path:
            self.merge_out_path.set(path)

    def _start_merge(self):
        files = list(self.merge_listbox.get(0, "end"))
        out_path = self.merge_out_path.get().strip()

        if len(files) < 2:
            messagebox.showerror("Error", "Add at least 2 .pptx files to merge.")
            return
        for f in files:
            if not os.path.isfile(f):
                messagebox.showerror("Error", f"File not found:\n{f}")
                return
        if not out_path:
            messagebox.showerror("Error", "Please choose a save location for the merged file.")
            return

        self.merge_btn.config(state="disabled")
        self.progress.start(10)
        self.status_var.set(f"Merging {len(files)} files…")

        threading.Thread(
            target=self._run_merge,
            args=(files, out_path),
            daemon=True
        ).start()

    def _run_merge(self, files: list, out_path: str):
        try:
            total = merge_pptx_files(files, out_path)
            self.root.after(0, self._on_merge_success, total, len(files), out_path)
        except Exception as exc:
            self.root.after(0, self._on_merge_error, str(exc))

    def _on_merge_success(self, slide_count: int, file_count: int, out_path: str):
        self.progress.stop()
        self.merge_btn.config(state="normal")
        self.status_var.set(f"Merged {file_count} files ({slide_count} slides) → {out_path}")
        messagebox.showinfo("Merge Complete",
            f"{file_count} files merged into {slide_count} slides.\nSaved to:\n{out_path}")

    def _on_merge_error(self, msg: str):
        self.progress.stop()
        self.merge_btn.config(state="normal")
        self.status_var.set(f"Merge error: {msg}")
        messagebox.showerror("Merge Failed", msg)

    # ------------------------------------------------------------------
    # Dependency check on startup
    # ------------------------------------------------------------------

    def _check_deps(self):
        missing = []
        if not PROTO_AVAILABLE:
            missing.append("• Proto bindings missing from proto_pb2/ folder")
        if not PPTX_AVAILABLE:
            missing.append("• python-pptx not installed  (pip install python-pptx)")
        if not STRIPRTF_AVAILABLE:
            missing.append("• striprtf not installed  (pip install striprtf)")
        if not WATCHDOG_AVAILABLE:
            missing.append("• watchdog not installed  (pip install watchdog) — needed for Watch Folder tab")

        if missing:
            msg = (
                "Some dependencies are missing:\n\n"
                + "\n".join(missing)
                + "\n\nSee README for setup instructions."
            )
            self.status_var.set("WARNING: missing dependencies — see startup dialog")
            messagebox.showwarning("Missing Dependencies", msg)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    root = tk.Tk()
    app = ConverterApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
